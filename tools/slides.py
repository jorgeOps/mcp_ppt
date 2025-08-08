"""
AutoSlides MCP – Herramientas de composición de diapositivas
===========================================================
Versión: 2025‑08‑07 b

Mejoras de layout:
* El área de viñetas ocupa ~55 % de ancho a la izquierda.
* Las imágenes se colocan en la derecha sin tapar texto.
* Márgenes consistentes y título siempre visible en la parte superior.
"""

from __future__ import annotations

import os
import time
from io import BytesIO
from pathlib import Path
from typing import List, Sequence

import httpx
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Configuración
# ---------------------------------------------------------------------------

OUT_DIR = Path("slides")
OUT_DIR.mkdir(exist_ok=True)

TEMPLATE_PATH = os.getenv("SLIDE_TEMPLATE")
IMAGE_MAX_W = Inches(3.5)  # algo más estrecho para caber en 45 % de ancho
IMAGE_MAX_H = Inches(3)
LEFT_COL_RATIO = 0.55       # 55 % viñetas, 45 % imágenes
MARGIN_X = Inches(0.4)
MARGIN_Y_TOP = Inches(1.5)  # deja sitio para título


# ---------------------------------------------------------------------------
# Presentación global
# ---------------------------------------------------------------------------

class _Deck:
    """Maneja una presentación única y helpers de layout."""

    def __init__(self):
        self.prs = Presentation(TEMPLATE_PATH) if TEMPLATE_PATH and Path(TEMPLATE_PATH).exists() else Presentation()
        self.width, self.height = self.prs.slide_width, self.prs.slide_height

    # ------------------------------------------------------------------
    # Añadir diapositivas
    # ------------------------------------------------------------------

    def add_bullet_slide(
        self,
        title: str,
        bullets: Sequence[str],
        images: Sequence[str] | None = None,
        notes: str | None = None,
    ) -> None:
        """Inserta diapositiva con título, viñetas a la izquierda e imágenes a la derecha."""

        layout = self._pick_layout()
        slide = self.prs.slides.add_slide(layout)

        # ---------- Título ----------
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # ---------- Placeholder cuerpo ----------
        body_shape = self._get_body_placeholder(slide)
        body_tf = body_shape.text_frame
        body_tf.clear()

        # Ajusta body a columna izquierda
        body_shape.left = MARGIN_X
        body_shape.top = MARGIN_Y_TOP
        body_shape.width = int(self.width * LEFT_COL_RATIO) - int(MARGIN_X * 2)
        body_shape.height = self.height - int(MARGIN_Y_TOP + MARGIN_X)

        for idx, bullet in enumerate(bullets):
            para = body_tf.add_paragraph() if idx else body_tf.paragraphs[0]
            para.text = bullet
            para.level = 0
            para.font.size = Pt(18)

        # ---------- Imágenes ----------
        if images:
            self._place_images(slide, images[:4])

        # ---------- Notas ----------
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # ------------------------------------------------------------------
    # Guardar
    # ------------------------------------------------------------------

    def save(self, filename: str | Path | None = None) -> str:
        if not filename:
            filename = f"deck_{time.strftime('%Y%m%d_%H%M%S')}.pptx"
        path = OUT_DIR / filename
        self.prs.save(path)
        return str(path.resolve())

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _pick_layout(self):
        layouts = self.prs.slide_layouts
        return layouts[1] if len(layouts) > 1 else layouts[0]

    @staticmethod
    def _get_body_placeholder(slide):
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                return shape
        return slide.shapes.placeholders[1]

    def _place_images(self, slide, urls: Sequence[str]):
        # Región derecha
        img_left = int(self.width * LEFT_COL_RATIO) + int(MARGIN_X)
        avail_w = self.width - img_left - int(MARGIN_X)
        avail_h = self.height - int(MARGIN_Y_TOP)

        cols = 1 if len(urls) == 1 else 2
        rows = (len(urls) + cols - 1) // cols

        img_w = min(IMAGE_MAX_W, avail_w / cols)
        img_h = min(IMAGE_MAX_H, avail_h / rows)

        col = row = 0
        top_start = MARGIN_Y_TOP

        for url in urls:
            try:
                data = httpx.get(url, timeout=10).content
                slide.shapes.add_picture(
                    BytesIO(data),
                    left=img_left + col * (img_w + MARGIN_X / 2),
                    top=top_start + row * (img_h + MARGIN_X / 2),
                    width=img_w,
                    height=img_h,
                )
            except Exception as exc:  # noqa: BLE001
                print(f"[WARN] No se pudo descargar imagen {url}: {exc}")

            col += 1
            if col >= cols:
                col = 0
                row += 1

# Instancia global
_deck = _Deck()

# ---------------------------------------------------------------------------
# API MCP
# ---------------------------------------------------------------------------

def create_slide(title: str, bullets: List[str], images: List[str] | None = None, notes: str | None = None) -> None:
    _deck.add_bullet_slide(title, bullets, images, notes)


def export_pptx(filename: str | None = None) -> str:
    return _deck.save(filename)
